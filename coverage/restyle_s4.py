"""Restyle slide 4 of the vRC deck to match slide 3's system."""
import os, sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.opc.package import Part
from pptx.opc.packuri import PackURI
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

SRC = sys.argv[1] if len(sys.argv) > 1 else "vrc.pptx"
OUT = sys.argv[2] if len(sys.argv) > 2 else "vrc_out.pptx"
ICONS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "icons")

NAVY = RGBColor(0x00, 0x28, 0x55)
BLACK = RGBColor(0x00, 0x00, 0x00)
BLUE = RGBColor(0x00, 0x67, 0xA5)      # slide 3 accent
HEAD, BODY = "Segoe UI Semibold", "Segoe UI"

SVG_URI = "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"
_seq = [900]


def drop(shape):
    shape._element.getparent().remove(shape._element)


def add_icon(slide, slug, x, y, size, name=None):
    png, svg = os.path.join(ICONS, slug + ".png"), os.path.join(ICONS, slug + ".svg")
    if not os.path.exists(png):
        return None
    pic = slide.shapes.add_picture(png, Inches(x), Inches(y), Inches(size), Inches(size))
    pic.name = "Icon %s" % (name or slug)
    if os.path.exists(svg):
        _seq[0] += 1
        with open(svg, "rb") as fh:
            blob = fh.read()
        part = Part(PackURI("/ppt/media/icon%d.svg" % _seq[0]), "image/svg+xml",
                    slide.part.package, blob)
        rId = slide.part.relate_to(part, RT.IMAGE)
        blip = pic._element.blipFill.find(qn("a:blip"))
        blip.append(parse_xml(
            '<a:extLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            '<a:ext uri="%s"><asvg:svgBlip'
            ' xmlns:asvg="http://schemas.microsoft.com/office/drawing/2016/SVG/main"'
            ' xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"'
            ' r:embed="%s"/></a:ext></a:extLst>' % (SVG_URI, rId)))
    return pic


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tf


def write(tf, runs, size=8.5, color=BLACK, font=BODY, align=PP_ALIGN.LEFT,
          line_spacing=1.06, space_after=0):
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    p.space_after = Pt(space_after)
    for text, opts in runs:
        r = p.add_run()
        r.text = text
        f = r.font
        f.name = opts.get("font", font)
        f.size = Pt(opts.get("size", size))
        f.bold = opts.get("bold", False)
        f.color.rgb = opts.get("color", color)
    return tf


def rule(slide, x, y, w, color=NAVY, pt=1.0):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y),
                                   Inches(x + w), Inches(y))
    c.line.color.rgb = color
    c.line.width = Pt(pt)
    c.shadow.inherit = False
    return c


def restyle_header(slide, tb, x, y, w, label):
    """Slide 3 treatment: black 10pt caps, thin navy rule beneath."""
    tb.left, tb.top, tb.width, tb.height = (Inches(x), Inches(y), Inches(w), Inches(0.17))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    for para in tf.paragraphs[1:]:
        para._p.getparent().remove(para._p)
    p = tf.paragraphs[0]
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    r = p.add_run()
    r.text = label
    r.font.name, r.font.size, r.font.bold = HEAD, Pt(10), True
    r.font.color.rgb = BLACK
    rule(slide, x, y + 0.20, w)


prs = Presentation(SRC)
s = list(prs.slides)[3]
shp = list(s.shapes)

# ------------------------------------------------------------------ deletions
for i in (24, 22, 21, 19, 7, 6, 5, 4):          # bands, takeaway, eyebrow, dd text
    drop(shp[i])

# ------------------------------------------------------------------- headers
restyle_header(s, shp[8], 0.25, 0.86, 5.45, "COMPETITIVE LANDSCAPE")
restyle_header(s, shp[20], 5.80, 0.86, 3.95, "DEMAND DRIVERS")
restyle_header(s, shp[23], 0.25, 5.26, 9.50, "INDICATIVE BUYER UNIVERSE")

# -------------------------------------------------- taller competitive quadrant
OLD_TOP, OLD_MID, OLD_BOT = 1.90, 2.83, 3.75
NEW_TOP, NEW_H = 1.18, 3.72
NEW_MID = NEW_TOP + NEW_H / 2

for idx, top, height in ((16, NEW_TOP, NEW_H), (15, NEW_TOP, NEW_H)):
    shp[idx].top, shp[idx].height = Inches(top), Inches(height)
shp[14].top = Inches(NEW_MID)                       # horizontal divider
for idx in (9, 44):                                 # highlighted top-right cell
    shp[idx].top, shp[idx].height = Inches(NEW_TOP), Inches(NEW_H / 2)

TITLES = {10: (3.17, 0), 11: (3.17, 1), 12: (0.56, 0), 13: (0.56, 1)}
for idx, (_, row) in TITLES.items():
    shp[idx].top = Inches((NEW_TOP if row == 0 else NEW_MID) + 0.09)

# regroup the logos: keep each cluster's internal arrangement, recentre in its cell
LOGOS = list(range(25, 44))
cells = {}
for idx in LOGOS:
    sh = shp[idx]
    col = 0 if sh.left < Inches(3.08) else 1
    row = 0 if sh.top < Inches(OLD_MID) else 1
    cells.setdefault((col, row), []).append(sh)

for (col, row), items in cells.items():
    lo = min(s_.top for s_ in items) / 914400
    hi = max((s_.top + s_.height) for s_ in items) / 914400
    cell_top = NEW_TOP if row == 0 else NEW_MID
    band_lo, band_hi = cell_top + 0.30, cell_top + NEW_H / 2 - 0.08
    shift = (band_lo + band_hi) / 2 - (lo + hi) / 2
    for s_ in items:
        s_.top = Inches(round(s_.top / 914400 + shift, 4))

shp[17].top = Inches(NEW_TOP + NEW_H + 0.06)        # x-axis caption
shp[18].top = Inches(NEW_MID - 0.06)                # y-axis caption (rotated)

# ------------------------------------------------- demand drivers, icon per row
DRIVERS = [
    ("dd_ai", "Data Quality Now Gates AI: ",
     "As companies adopt more AI systems, input data quality is increasingly "
     "important for output trust and reliability"),
    ("dd_reg", "Regulatory and Cost Pressure: ",
     "Regulators are demanding higher visibility into where data was originated, "
     "how it was used and more"),
    ("dd_volume", "Data Volume: ",
     "As the economy digitizes further there is more data and data is more complex; "
     "streamlining that data becomes more valuable"),
    ("dd_cloud", "Banks are Modernizing: ",
     "Moving from legacy on-premise software onto cloud solutions"),
    ("dd_blindspot", "Limited Visibility: ",
     "With little insight into data sets, companies unknowingly pay for duplicate "
     "vendors, handle data manually and read inconsistent data"),
]
DX, DW = 5.80, 3.95
ICON = 0.24
rh = NEW_H / len(DRIVERS)
for i, (slug, leadin, rest) in enumerate(DRIVERS):
    ry = NEW_TOP + i * rh
    add_icon(s, slug, DX, ry + 0.06, ICON, leadin.strip(": "))
    tf = textbox(s, DX + ICON + 0.14, ry, DW - ICON - 0.14, rh - 0.10,
                 anchor=MSO_ANCHOR.TOP)
    write(tf, [(leadin, {"bold": True, "font": HEAD, "color": NAVY}),
               (rest, {})], size=8.5)

prs.save(OUT)
print("wrote", OUT)
