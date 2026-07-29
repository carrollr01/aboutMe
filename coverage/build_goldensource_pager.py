"""Build the GoldenSource coverage two-pager on the Houlihan Lokey template.

Usage:  python build_goldensource_pager.py [path/to/HL_template.pptx] [out.pptx]

The HL template is NOT committed to this repo (it carries the firm's
"Strictly Confidential. Not for Distribution." cover). Point the first argument
at a local copy of the HL Refresh 2023 deck; the script inherits its master,
theme, fonts and footer logo, and writes only the two content pages.

Charts are native PowerPoint charts with embedded Excel workbooks; diagrams are
native shapes. Nothing on either page is a picture.
"""

import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_TICK_MARK
from pptx.oxml.ns import qn

SRC = sys.argv[1] if len(sys.argv) > 1 else "source.pptx"
OUT = sys.argv[2] if len(sys.argv) > 2 else "GoldenSource_Coverage_Twopager.pptx"

# ---------------------------------------------------------------- HL palette
NAVY = RGBColor(0x00, 0x28, 0x55)   # theme dk2 - primary for main slides
MID = RGBColor(0x50, 0x8B, 0xC9)    # theme accent1
GRAY = RGBColor(0x52, 0x57, 0x66)   # theme dk1 - body copy
SLATE = RGBColor(0x7E, 0x85, 0x97)  # theme accent3
LGRAY = RGBColor(0xBC, 0xBF, 0xC6)  # theme accent2 - rules
TEAL = RGBColor(0x24, 0xB1, 0xB1)   # theme accent4 - reserved for punchlines
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ROW = RGBColor(0xEE, 0xF2, 0xF7)    # table row striping only

# chevron progression, navy -> theme accent1
FLOW = [RGBColor(0x00, 0x28, 0x55), RGBColor(0x1C, 0x44, 0x70),
        RGBColor(0x38, 0x61, 0x8B), RGBColor(0x50, 0x8B, 0xC9)]

HEAD = "Segoe UI Semibold"
BODY = "Segoe UI"

# ------------------------------------------------------------------- geometry
LX = 0.25
FULL_W = 9.50
BAND_Y, BAND_H = 0.90, 0.58
BOT = 6.95
SEC_H = 0.24          # slim section header band
PAD = 0.10


# --------------------------------------------------------------------- helpers
def delete_all_slides(prs):
    lst = prs.slides._sldIdLst
    for sld in list(lst):
        prs.part.drop_rel(sld.get(qn("r:id")))
        lst.remove(sld)


def no_line(shape):
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def rect(slide, x, y, w, h, fill=None, shape=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    return no_line(sh)


def hairline(slide, x, y, w, color=LGRAY, weight=0.008):
    return rect(slide, x, y, w, weight, color)


def vrule(slide, x, y, h, color=LGRAY, weight=0.008):
    return rect(slide, x, y, weight, h, color)


def set_bullet(para, color=MID, char="▪", font="Arial", size_pct=90):
    """Attach a square bullet via pPr - never a literal character in the text."""
    pPr = para._p.get_or_add_pPr()
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum", "a:buClr", "a:buFont", "a:buSzPct"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    buClr = pPr.makeelement(qn("a:buClr"), {})
    buClr.append(pPr.makeelement(qn("a:srgbClr"), {"val": str(color)}))
    for el in (buClr,
               pPr.makeelement(qn("a:buSzPct"), {"val": str(size_pct * 1000)}),
               pPr.makeelement(qn("a:buFont"), {"typeface": font}),
               pPr.makeelement(qn("a:buChar"), {"char": char})):
        pPr.append(el)


def no_bullet(para):
    pPr = para._p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum", "a:buClr", "a:buFont", "a:buSzPct"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    if not pPr.findall(qn("a:buNone")):
        pPr.append(pPr.makeelement(qn("a:buNone"), {}))


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tf


def write(tf, paras, size=8.0, color=GRAY, font=BODY, align=PP_ALIGN.LEFT,
          space_after=4, line_spacing=1.02):
    for i, spec in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = spec.get("align", align)
        p.line_spacing = spec.get("line_spacing", line_spacing)
        p.space_after = Pt(spec.get("space_after", space_after))
        p.space_before = Pt(spec.get("space_before", 0))
        pPr = p._p.get_or_add_pPr()
        if spec.get("bullet"):
            pPr.set("marL", str(int(Inches(0.115))))
            pPr.set("indent", str(int(-Inches(0.115))))
            set_bullet(p, color=spec.get("bullet_color", MID))
        else:
            pPr.set("marL", "0")
            pPr.set("indent", "0")
            no_bullet(p)
        for text, opts in spec["runs"]:
            r = p.add_run()
            r.text = text
            f = r.font
            f.name = opts.get("font", font)
            f.size = Pt(opts.get("size", size))
            f.bold = opts.get("bold", False)
            f.italic = opts.get("italic", False)
            f.color.rgb = opts.get("color", color)
    return tf


def plain(text, **opts):
    return {"runs": [(text, opts)]}


def lead(label, rest, **opts):
    """A bold navy lead-in followed by body copy, in one paragraph."""
    return {"runs": [(label, {"bold": True, "color": NAVY, "font": HEAD}), (rest, {})], **opts}


def section(slide, x, y, w, label):
    """Slim navy header band. Returns the y at which content should start."""
    rect(slide, x, y, w, SEC_H, NAVY)
    tf = textbox(slide, x + PAD, y, w - 2 * PAD, SEC_H, anchor=MSO_ANCHOR.MIDDLE)
    write(tf, [plain(label, bold=True, size=8.5, color=WHITE, font=HEAD)], space_after=0)
    return y + SEC_H + 0.10


def footnote(slide, text):
    tf = textbox(slide, 2.05, 7.10, 7.70, 0.38)
    write(tf, [plain(text, size=7, color=SLATE, italic=True)],
          space_after=0, line_spacing=1.0)


def page_frame(prs, layouts, title, eyebrow, takeaway):
    slide = prs.slides.add_slide(layouts["Title Only"])
    t = slide.shapes.title
    t.text_frame.word_wrap = False
    r = t.text_frame.paragraphs[0].add_run()
    r.text = title
    r.font.name, r.font.size, r.font.bold = HEAD, Pt(18), True
    r.font.color.rgb = NAVY

    tf = textbox(slide, 4.60, 0.40, 5.15, 0.30, anchor=MSO_ANCHOR.BOTTOM)
    write(tf, [plain(eyebrow, size=8.5, color=SLATE, italic=True)],
          align=PP_ALIGN.RIGHT, space_after=0)

    rect(slide, LX, BAND_Y, FULL_W, BAND_H, NAVY)
    tf = textbox(slide, LX + 0.14, BAND_Y, FULL_W - 0.28, BAND_H, anchor=MSO_ANCHOR.MIDDLE)
    write(tf, [plain(takeaway, size=9.5, color=WHITE, font=HEAD, bold=True)],
          space_after=0, line_spacing=1.06)
    return slide


# ------------------------------------------------------------------- zone art
def kpi_strip(slide, x, y, w, h, items):
    """Open stat row: navy rule on top, hairline dividers, no fills."""
    hairline(slide, x, y, w, NAVY, 0.014)
    cw = w / len(items)
    for i, (num, label) in enumerate(items):
        cx = x + i * cw
        if i:
            vrule(slide, cx, y + 0.08, h - 0.10)
        tf = textbox(slide, cx + 0.14, y + 0.07, cw - 0.28, 0.24, anchor=MSO_ANCHOR.TOP)
        write(tf, [plain(num, size=11.5, color=NAVY, bold=True, font=HEAD)], space_after=0)
        tf = textbox(slide, cx + 0.14, y + 0.31, cw - 0.28, h - 0.33)
        write(tf, [plain(label, size=7.5, color=GRAY)], space_after=0, line_spacing=1.0)


def kv_rows(slide, x, y, w, rows, label_w=0.95, size=7.5):
    """Striped key/value rows. rows: (label, value, n_lines)."""
    cy = y
    for i, (label, value, nlines) in enumerate(rows):
        h = 0.205 if nlines == 1 else 0.32
        if i % 2 == 0:
            rect(slide, x - 0.06, cy, w + 0.12, h, ROW)
        tf = textbox(slide, x, cy, label_w, h, anchor=MSO_ANCHOR.MIDDLE)
        write(tf, [plain(label, size=size, color=NAVY, bold=True, font=HEAD)], space_after=0)
        tf = textbox(slide, x + label_w, cy, w - label_w, h, anchor=MSO_ANCHOR.MIDDLE)
        write(tf, [plain(value, size=size, color=GRAY)], space_after=0, line_spacing=1.0)
        cy += h
    return cy


def flow_diagram(slide, x, y, w, h, stages, chev_h=0.34):
    """Chevron process flow with a module list under each stage."""
    n = len(stages)
    overlap = 0.10
    cw = (w + overlap * (n - 1)) / n
    for i, (title, modules) in enumerate(stages):
        cx = x + i * (cw - overlap)
        ch = rect(slide, cx, y, cw, chev_h, FLOW[i % len(FLOW)], MSO_SHAPE.CHEVRON)
        tf = ch.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.16)
        tf.margin_right = Inches(0.06)
        tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(tf, [plain(title, size=7.5, color=WHITE, bold=True, font=HEAD)],
              align=PP_ALIGN.CENTER, space_after=0)
        tf = textbox(slide, cx + 0.18, y + chev_h + 0.07, cw - 0.30, h - chev_h - 0.07)
        write(tf, [{"runs": [(m, {})], "space_after": 1.5, "bullet": True,
                    "bullet_color": FLOW[i % len(FLOW)]} for m in modules],
              size=7, line_spacing=1.0)


def timeline(slide, x, y, w, h, nodes, label_size=7, dot=0.10):
    """Horizontal spine with year nodes; each node captions below the line."""
    n = len(nodes)
    cw = w / n
    axis_y = y + 0.20
    hairline(slide, x + cw * 0.5, axis_y, w - cw, NAVY, 0.011)
    for i, node in enumerate(nodes):
        year, lines, accent = node
        cx = x + i * cw
        mid = cx + cw / 2
        color = TEAL if accent else NAVY
        d = rect(slide, mid - dot / 2, axis_y - dot / 2 + 0.005, dot, dot, color, MSO_SHAPE.OVAL)
        d.line.fill.solid()
        d.line.fill.fore_color.rgb = WHITE
        d.line.width = Pt(1)
        tf = textbox(slide, cx + 0.04, y - 0.02, cw - 0.08, 0.18, anchor=MSO_ANCHOR.BOTTOM)
        write(tf, [plain(year, size=8, color=color, bold=True, font=HEAD)],
              align=PP_ALIGN.CENTER, space_after=0)
        tf = textbox(slide, cx + 0.04, axis_y + 0.13, cw - 0.08, h - (axis_y - y) - 0.13)
        write(tf, [{"runs": [(t, {"bold": b, "color": color if b else GRAY,
                                  "font": HEAD if b else BODY})], "space_after": 0.5}
                   for t, b in lines],
              size=label_size, align=PP_ALIGN.CENTER, line_spacing=1.0)


def outline(slide, x, y, w, h, color=LGRAY, weight=0.75):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.background()
    sh.line.color.rgb = color
    sh.line.width = Pt(weight)
    sh.shadow.inherit = False
    return sh


def matrix_2x2(slide, x, y, w, h, quadrants, x_axis, y_axis,
               highlight=(1, 0), pad=0.09, axis_w=0.22, axis_h=0.20):
    """Qualitative 2x2. quadrants keyed (col, row) with row 0 = top.
    Each value is (quadrant_title, [(text, opts), ...])."""
    gx, gy = x + axis_w, y
    gw, gh = w - axis_w, h - axis_h
    cw, ch = gw / 2, gh / 2

    for (col, row), (title, runs) in quadrants.items():
        qx, qy = gx + col * cw, gy + row * ch
        if (col, row) == highlight:
            rect(slide, qx, qy, cw, ch, ROW)
        tf = textbox(slide, qx + pad, qy + pad, cw - 2 * pad, 0.20)
        write(tf, [plain(title, size=7, color=NAVY, bold=True, font=HEAD)],
              space_after=0, line_spacing=1.0)
        tf = textbox(slide, qx + pad, qy + pad + 0.22, cw - 2 * pad, ch - pad * 2 - 0.22)
        write(tf, [{"runs": runs}], size=7, space_after=0, line_spacing=1.12)

    hairline(slide, gx, gy + ch, gw, LGRAY)
    vrule(slide, gx + cw, gy, gh, LGRAY)
    outline(slide, gx, gy, gw, gh)

    tf = textbox(slide, gx, gy + gh + 0.04, gw, axis_h - 0.04)
    write(tf, [plain(x_axis, size=7, color=SLATE, bold=True, font=HEAD)],
          align=PP_ALIGN.CENTER, space_after=0)
    tb = slide.shapes.add_textbox(Inches(x + axis_w / 2 - gh / 2), Inches(gy + gh / 2 - 0.10),
                                  Inches(gh), Inches(0.20))
    tb.rotation = 270
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    write(tf, [plain(y_axis, size=7, color=SLATE, bold=True, font=HEAD)],
          align=PP_ALIGN.CENTER, space_after=0)


def lanes(slide, x, y, w, h, rows, chip_w=1.24, gap=0.08):
    """Colour-chipped lanes: a category chip on the left, names on the right."""
    lh = (h - gap * (len(rows) - 1)) / len(rows)
    for i, (label, names) in enumerate(rows):
        ly = y + i * (lh + gap)
        chip = rect(slide, x, ly, chip_w, lh, FLOW[i % len(FLOW)])
        tf = chip.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.05)
        tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(tf, [plain(label, size=7, color=WHITE, bold=True, font=HEAD)],
              align=PP_ALIGN.CENTER, space_after=0, line_spacing=1.0)
        tf = textbox(slide, x + chip_w + 0.10, ly, w - chip_w - 0.10, lh,
                     anchor=MSO_ANCHOR.MIDDLE)
        write(tf, [plain(names, size=7.5, color=GRAY)], space_after=0, line_spacing=1.05)


def style_axes(chart, cat_size=7.5):
    chart.has_title = False
    chart.has_legend = False
    va = chart.value_axis
    va.visible = False
    va.has_major_gridlines = False
    ca = chart.category_axis
    ca.has_major_gridlines = False
    ca.major_tick_mark = XL_TICK_MARK.NONE
    ca.format.line.color.rgb = LGRAY
    ca.tick_labels.font.size = Pt(cat_size)
    ca.tick_labels.font.color.rgb = GRAY
    ca.tick_labels.font.name = BODY


def native_chart(slide, kind, x, y, w, h, categories, values, colors,
                 number_format="General", label_pos=XL_LABEL_POSITION.OUTSIDE_END,
                 gap_width=70, cat_size=7.5):
    cd = CategoryChartData()
    cd.categories = categories
    cd.add_series("Series 1", values)
    gf = slide.shapes.add_chart(kind, Inches(x), Inches(y), Inches(w), Inches(h), cd)
    chart = gf.chart
    style_axes(chart, cat_size)
    plot = chart.plots[0]
    plot.gap_width = gap_width
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.size = Pt(8)
    dl.font.bold = True
    dl.font.name = HEAD
    dl.font.color.rgb = NAVY
    dl.number_format = number_format
    dl.number_format_is_linked = False
    dl.position = label_pos
    ser = plot.series[0]
    for i, c in enumerate(colors):
        pt = ser.points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = c
    return chart


# ------------------------------------------------- platform pipeline visual
ENDS_BODY = RGBColor(0xF1, 0xF4, 0xF8)   # sources / consumers: outside the platform
STAGE = [RGBColor(0x00, 0x28, 0x55), RGBColor(0x1C, 0x44, 0x70), RGBColor(0x38, 0x61, 0x8B)]


def pipe_block(slide, x, y, w, h, title, items, head_fill, body_fill,
               body_color, bullet_color, head_h=0.24, outlined=False):
    rect(slide, x, y, w, h, body_fill)
    if outlined:
        outline(slide, x, y, w, h, LGRAY)
    rect(slide, x, y, w, head_h, head_fill)
    tf = textbox(slide, x + 0.05, y, w - 0.10, head_h, anchor=MSO_ANCHOR.MIDDLE)
    write(tf, [plain(title, size=7, color=WHITE, bold=True, font=HEAD)],
          align=PP_ALIGN.CENTER, space_after=0, line_spacing=1.0)
    tf = textbox(slide, x + 0.09, y + head_h + 0.09, w - 0.18, h - head_h - 0.14)
    write(tf, [{"runs": [(i, {})], "bullet": True, "bullet_color": bullet_color,
                "space_after": 2.0} for i in items],
          size=7, color=body_color, line_spacing=1.04)


def wire(slide, x0, y0, x1, y1, color=LGRAY, pt=0.75):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x0), Inches(y0),
                                   Inches(x1), Inches(y1))
    c.line.color.rgb = color
    c.line.width = Pt(pt)
    c.shadow.inherit = False
    return c


def node(slide, cx, cy, d, fill, ring=WHITE):
    sh = rect(slide, cx - d / 2, cy - d / 2, d, d, fill, MSO_SHAPE.OVAL)
    sh.line.fill.solid()
    sh.line.fill.fore_color.rgb = ring
    sh.line.width = Pt(1.25)
    return sh


def endpoints(slide, x, y, w, h, items, tie_x, tie_y, inbound=True, size=7):
    """Stacked labels, each wired to a single bundling point."""
    n = len(items)
    rh = h / n
    for i, t in enumerate(items):
        cy = y + i * rh + rh / 2
        tf = textbox(slide, x, cy - 0.10, w, 0.20, anchor=MSO_ANCHOR.MIDDLE)
        write(tf, [plain(t, size=size, color=GRAY)],
              align=PP_ALIGN.RIGHT if inbound else PP_ALIGN.LEFT,
              space_after=0, line_spacing=1.0)
        if inbound:
            wire(slide, x + w + 0.07, cy, tie_x, tie_y)
        else:
            wire(slide, tie_x, tie_y, x - 0.07, cy)


def marquee(slide, x, y, w, h, rows):
    rh = h / len(rows)
    for i, (dot, name, desc) in enumerate(rows):
        ry = y + i * rh
        rect(slide, x, ry + 0.045, 0.08, 0.08, dot, MSO_SHAPE.OVAL)
        tf = textbox(slide, x + 0.145, ry, w - 0.145, 0.17)
        write(tf, [plain(name, size=8, color=NAVY, bold=True, font=HEAD)], space_after=0)
        tf = textbox(slide, x + 0.145, ry + 0.165, w - 0.145, rh - 0.185)
        write(tf, [plain(desc, size=7, color=GRAY)], space_after=0, line_spacing=1.0)


def platform_panel(slide, x, y, w, h):
    """9.30in x 2.20in: marquee products left, bundled wire flow right."""
    l_w, gut = 2.50, 0.20
    r_x, r_w = x + l_w + gut, w - l_w - gut

    marquee(slide, x, y, l_w, h, [
        (STAGE[1], "Securities Master", "Cross-asset instrument mastering, ESG-enabled"),
        (STAGE[1], "Price Master", "Rule-based multi-vendor pricing and valuation"),
        (STAGE[2], "OMNI", "Snowflake Native App; runs in the client\u2019s own account"),
        (STAGE[2], "IBOR", "Real-time investment book of record"),
        (TEAL, "Scout", "AI layer on Amazon Bedrock, June 2026"),
    ])

    src_w, con_w, lead = 0.88, 1.05, 0.40
    bundle = r_x + src_w + lead
    unbundle = r_x + r_w - con_w - lead
    trunk_y = y + 1.02
    stack_h = 1.70
    stack_y = trunk_y - stack_h / 2

    endpoints(slide, r_x, stack_y, src_w, stack_h,
              ["Data vendors", "Exchanges", "Index providers", "Custodians",
               "Fund admins", "Internal systems", "Counterparties"],
              bundle, trunk_y, inbound=True)
    endpoints(slide, r_x + r_w - con_w, stack_y, con_w, stack_h,
              ["Trading and OMS", "Risk and capital", "Regulatory reporting",
               "Client reporting", "Fund accounting", "AI and analytics"],
              unbundle, trunk_y, inbound=False)

    # the trunk: bundled, verified, one record
    rect(slide, bundle, trunk_y - 0.011, unbundle - bundle, 0.022, NAVY)

    span = unbundle - bundle
    stages = [
        (0.16, 0.15, "CONNECT", "Adaptors and 100+ pre-built vendor feeds"),
        (0.50, 0.19, "MASTER",
         "One verified record: securities, entity, price, corporate actions, ESG"),
        (0.84, 0.15, "DISTRIBUTE", "Data Warehouse, OMNI into Snowflake, IBOR and APIs"),
    ]
    cap_w = 1.24
    for frac, d, ttl, cap in stages:
        cx = bundle + span * frac
        node(slide, cx, trunk_y, d, NAVY)
        tf = textbox(slide, cx - cap_w / 2, trunk_y + 0.14, cap_w, 0.16)
        write(tf, [plain(ttl, size=7, color=NAVY, bold=True, font=HEAD)],
              align=PP_ALIGN.CENTER, space_after=0)
        tf = textbox(slide, cx - cap_w / 2, trunk_y + 0.31, cap_w, 0.56)
        write(tf, [plain(cap, size=7, color=GRAY)],
              align=PP_ALIGN.CENTER, space_after=0, line_spacing=1.04)

    # Scout taps the trunk from above
    scout_cx = bundle + span * 0.50
    pill_w, pill_h = 2.30, 0.24
    py_ = y + 0.10
    sp_ = rect(slide, scout_cx - pill_w / 2, py_, pill_w, pill_h, WHITE,
               MSO_SHAPE.ROUNDED_RECTANGLE)
    sp_.line.fill.solid()
    sp_.line.fill.fore_color.rgb = TEAL
    sp_.line.width = Pt(1)
    tf = textbox(slide, scout_cx - pill_w / 2 + 0.08, py_, pill_w - 0.16, pill_h,
                 anchor=MSO_ANCHOR.MIDDLE)
    write(tf, [plain("SCOUT   \u00b7   AI reasoning layer on Amazon Bedrock",
                     size=7, color=TEAL, bold=True, font=HEAD)],
          align=PP_ALIGN.CENTER, space_after=0)
    wire(slide, scout_cx, py_ + pill_h, scout_cx, trunk_y - 0.10, TEAL, 1)

    br = y + h - 0.20
    hairline(slide, bundle, br, span, NAVY, 0.011)
    tf = textbox(slide, bundle, br + 0.025, span, 0.15)
    write(tf, [plain("GOLDENSOURCE PLATFORM", size=7, color=NAVY, bold=True, font=HEAD)],
          align=PP_ALIGN.CENTER, space_after=0)


# ================================================================ build deck
prs = Presentation(SRC)
delete_all_slides(prs)
layouts = {l.name: l for l in prs.slide_master.slide_layouts}

# ==================================================================== PAGE 1
s1 = page_frame(
    prs, layouts,
    "GoldenSource",
    "Coverage Profile  |  Data & Analytics  |  July 2026",
    "GoldenSource has mastered securities, pricing and entity data for capital markets since "
    "1984. Gemspring Capital has owned it since May 2022; headcount is up roughly 60% over that "
    "period and 2025 was a record year.",
)

# --- zone A: open KPI strip, full width
kpi_strip(s1, LX, 1.56, FULL_W, 0.56, [
    ("1984", "Founded; over 40 years in capital markets data"),
    ("100+", "Pre-built vendor data feeds and adaptors"),
    ("6", "Offices across North America, EMEA and APAC"),
    ("4 yrs", "Gemspring Capital hold period to date"),
])

# --- zone B: overview (wide left) + snapshot (narrow right)
cy = section(s1, LX, 2.24, 5.45, "COMPANY OVERVIEW")
tf = textbox(s1, LX + 0.02, cy, 5.41, 4.40 - cy)
write(tf, [
    {"bullet": True, "runs": [
        ("Founded in 1984 and headquartered in New York, GoldenSource is one of the original "
         "enterprise data management (EDM) vendors in capital markets and is credited with "
         "popularising the term.", {})]},
    {"bullet": True, "runs": [
        ("Masters, governs and distributes securities, entity, counterparty, client, pricing, "
         "position, corporate-action and ESG data for banks, asset managers and insurers.", {})]},
    {"bullet": True, "runs": [
        ("One platform, two go-to-market lenses: the ", {}),
        ("Investment Data Platform", {"bold": True, "color": NAVY, "font": HEAD}),
        (" for the buy side and a ", {}),
        ("Trading, Risk and Regulatory Data Platform", {"bold": True, "color": NAVY, "font": HEAD}),
        (" for the sell side.", {})]},
    {"bullet": True, "runs": [
        ("Recurring subscription licence plus cloud hosting and managed services, deployed SaaS "
         "or on-premise; founding sponsor of the EDM Council.", {})]},
    {"bullet": True, "runs": [
        ("Leadership refreshed under sponsor ownership: James Corrigan became CEO in September "
         "2024, succeeding John Eley after a decade, with a Chief Customer Officer seat created "
         "in 2026.", {})]},
])

cy = section(s1, 5.80, 2.24, 3.95, "COMPANY SNAPSHOT")
kv_rows(s1, 5.86, cy - 0.03, 3.83, [
    ("Founded", "1984  |  New York, NY", 1),
    ("Ownership", "Gemspring Capital (May 2022); previously The Invus Group", 2),
    ("CEO", "James Corrigan (since Sept. 2024)", 1),
    ("Employees", "~620 est., versus 380 at the 2022 close", 1),
    ("Revenue", "~$63M est. (third-party; not company-reported)", 2),
    ("Offices", "New York, London, Milan, Mumbai, Melbourne, Hong Kong", 2),
    ("Delivery", "SaaS on AWS, Azure and GCP; AWS Marketplace; on-premise", 2),
])

# --- zone C: compact platform flow (the full 9.30 x 2.20 panel ships separately)
cy = section(s1, LX, 4.58, FULL_W, "PLATFORM AND PRODUCT SUITE")
flow_diagram(s1, LX, cy, FULL_W, 5.76 - cy, [
    ("SOURCE AND CONNECT", ["Connections and Adaptors", "100+ vendor feeds"]),
    ("MASTER AND GOVERN", ["Securities and Entity Master", "Price Master, corporate actions, ESG"]),
    ("DISTRIBUTE", ["Data Warehouse", "OMNI (Snowflake Native App)", "Real-time IBOR"]),
    ("CONSUME AND REASON", ["Scout on Amazon Bedrock", "Chat plus MCP agent builder"]),
])

# --- zone D: milestone timeline (wide) + native headcount chart (narrow)
cy = section(s1, LX, 5.86, 6.05, "SELECTED MILESTONES")
timeline(s1, LX, cy, 6.05, BOT - cy, [
    ("1984", [("Founded in", False), ("New York", False)], False),
    ("2022", [("Gemspring buys", False), ("from Invus", False)], True),
    ("2023\u201324", [("IBOR, OMNI and", False), ("V10; new CEO", False)], False),
    ("2025", [("Record year;", False), ("flagship wins", False)], False),
    ("2026", [("Scout AI", False), ("launches", False)], False),
])

cy = section(s1, 6.55, 5.86, 3.20, "HEADCOUNT UNDER GEMSPRING")
native_chart(s1, XL_CHART_TYPE.COLUMN_CLUSTERED, 6.52, cy - 0.09, 3.26, BOT - cy + 0.05,
             ["2022", "2026E"], (380, 620), [SLATE, NAVY], gap_width=60)

footnote(s1, "Sources: GoldenSource and Gemspring Capital press releases; Businesswire; Finextra; "
             "WatersTechnology; company website. Employee and revenue figures are third-party "
             "estimates and are not company-reported; confirm in diligence.")

# ============================================================ PAGE: PLATFORM
# 9.3in of usable width: marquee products on the left, data pipeline on the right.
PW = 9.30
PX = 0.35
L_W = 2.45          # left panel: marquee products
GUT = 0.25
R_X = PX + L_W + GUT
R_W = PW - L_W - GUT

STAGE = [RGBColor(0x00, 0x28, 0x55), RGBColor(0x1C, 0x44, 0x70), RGBColor(0x38, 0x61, 0x8B)]
ENDS = RGBColor(0xF1, 0xF4, 0xF8)   # sources / consumers, outside the platform


def mini_label(slide, x, y, w, text):
    tf = textbox(slide, x, y, w, 0.18, anchor=MSO_ANCHOR.BOTTOM)
    write(tf, [plain(text, size=7.5, color=SLATE, bold=True, font=HEAD)], space_after=0)
    return y + 0.24


def product_card(slide, x, y, w, h, dot, name, desc, rule=True):
    d = rect(slide, x, y + 0.055, 0.085, 0.085, dot, MSO_SHAPE.OVAL)
    tf = textbox(slide, x + 0.16, y, w - 0.16, 0.20)
    write(tf, [plain(name, size=9, color=NAVY, bold=True, font=HEAD)], space_after=0)
    tf = textbox(slide, x + 0.16, y + 0.21, w - 0.16, h - 0.26)
    write(tf, [plain(desc, size=7.5, color=GRAY)], space_after=0, line_spacing=1.1)
    if rule:
        hairline(slide, x, y + h + 0.015, w, LGRAY)


def stage_block(slide, x, y, w, h, title, items, fill, text_on_fill=True,
                head_h=0.34, border=False):
    if border:
        outline(slide, x, y, w, h, LGRAY)
        rect(slide, x, y, w, head_h, fill)
    else:
        rect(slide, x, y, w, h, fill)
    tf = textbox(slide, x + 0.08, y, w - 0.16, head_h, anchor=MSO_ANCHOR.MIDDLE)
    write(tf, [plain(title, size=7.5, color=WHITE if text_on_fill else NAVY,
                     bold=True, font=HEAD)],
          align=PP_ALIGN.CENTER, space_after=0, line_spacing=1.0)
    body = WHITE if text_on_fill else GRAY
    tf = textbox(slide, x + 0.13, y + head_h + 0.12, w - 0.26, h - head_h - 0.20)
    write(tf, [{"runs": [(i, {})], "bullet": True,
                "bullet_color": WHITE if text_on_fill else MID, "space_after": 3.5}
               for i in items],
          size=7, color=body, line_spacing=1.08)


def flow_arrow(slide, x, y, size, color=LGRAY):
    return rect(slide, x, y, size, size, color, MSO_SHAPE.RIGHT_ARROW)


sp = page_frame(
    prs, layouts,
    "GoldenSource",
    "Platform and Product Suite  |  July 2026",
    "One platform sits between the client’s data vendors and every system that "
    "consumes the data. GoldenSource connects the sources, masters and governs the record, "
    "and distributes it to trading, risk, reporting and, since June 2026, to AI agents.",
)

cy = section(sp, PX, 1.58, PW, "PLATFORM AND PRODUCT SUITE")

# ---- left panel: marquee products
ly = mini_label(sp, PX, cy, L_W, "MARQUEE PRODUCTS")
cards = [
    (STAGE[1], "Securities Master",
     "Cross-asset instrument mastering across equities, fixed income, funds, "
     "structured products and derivatives. ESG-enabled."),
    (STAGE[1], "Price Master",
     "Rule-based multi-vendor pricing and valuation, centralised for NAV and "
     "portfolio valuation."),
    (STAGE[2], "OMNI",
     "Snowflake Native App. Runs the GoldenSource data model inside the client’s "
     "own Snowflake account; data never leaves."),
    (STAGE[2], "IBOR",
     "Real-time investment book of record. Positions and valuations update "
     "intraday as prices and trades arrive."),
    (TEAL, "Scout",
     "AI layer launched June 2026. Chat interface and MCP agent builder over the "
     "mastered record, on Amazon Bedrock."),
]
card_h, card_gap = 0.775, 0.07
for i, (dot, name, desc) in enumerate(cards):
    product_card(sp, PX, ly + i * (card_h + card_gap), L_W, card_h, dot, name, desc,
                 rule=(i < len(cards) - 1))

# ---- right panel: the pipeline
ry = mini_label(sp, R_X, cy, R_W, "DATA PIPELINE")

ARR = 0.24
src_w, con_w = 1.10, 1.02
plat_w = R_W - src_w - con_w - 2 * ARR
c_w, m_w, d_w = plat_w * 0.30, plat_w * 0.40, plat_w * 0.30

scout_x = R_X + src_w + ARR
scout_h = 0.34
rect(sp, scout_x, ry, plat_w, scout_h, TEAL)
tf = textbox(sp, scout_x + 0.10, ry, plat_w - 0.20, scout_h, anchor=MSO_ANCHOR.MIDDLE)
write(tf, [plain("SCOUT  ·  AI reasoning layer over the mastered record, on Amazon Bedrock",
                 size=7.5, color=WHITE, bold=True, font=HEAD)],
      align=PP_ALIGN.CENTER, space_after=0)

py_ = ry + scout_h + 0.13
ph = 6.30 - py_

stage_block(sp, R_X, py_, src_w, ph, "SOURCES", [
    "Market data vendors", "Exchanges and index providers", "Custodians and administrators",
    "Internal trading and accounting systems", "Client and counterparty records",
], ENDS, text_on_fill=False, border=True)
flow_arrow(sp, R_X + src_w + 0.02, py_ + ph / 2 - ARR / 2, ARR - 0.04)

cx = scout_x
stage_block(sp, cx, py_, c_w, ph, "CONNECT", [
    "Connections and Adaptors", "100+ pre-built vendor feeds",
    "Validation and exception rules",
], STAGE[0])
cx += c_w
stage_block(sp, cx, py_, m_w, ph, "MASTER AND GOVERN", [
    "Securities Master", "Entity and Customer Master", "Price Master and valuations",
    "Corporate actions", "ESG Impact Plus", "Lineage, audit and entitlements",
], STAGE[1])
cx += m_w
stage_block(sp, cx, py_, d_w, ph, "DISTRIBUTE", [
    "Data Warehouse", "OMNI into Snowflake", "Real-time IBOR", "APIs and downstream feeds",
], STAGE[2])

flow_arrow(sp, scout_x + plat_w + 0.02, py_ + ph / 2 - ARR / 2, ARR - 0.04)
stage_block(sp, R_X + R_W - con_w, py_, con_w, ph, "CONSUMERS", [
    "Trading and order management", "Risk and capital", "Regulatory reporting",
    "Client and fund reporting", "Analytics and AI models",
], ENDS, text_on_fill=False, border=True)

# platform bracket
hairline(sp, scout_x, 6.42, plat_w, NAVY, 0.011)
tf = textbox(sp, scout_x, 6.46, plat_w, 0.18)
write(tf, [plain("GOLDENSOURCE PLATFORM", size=7.5, color=NAVY, bold=True, font=HEAD)],
      align=PP_ALIGN.CENTER, space_after=0)

tf = textbox(sp, PX, 6.68, PW, 0.24)
write(tf, [plain("Deployed as SaaS on AWS, Azure or GCP, inside the client’s own Snowflake "
                 "account via OMNI, or on-premise. Modules are separately licensable; three "
                 "are listed standalone on AWS Marketplace.",
                 size=7.5, color=GRAY)], space_after=0, line_spacing=1.05)

footnote(sp, "Sources: GoldenSource product documentation and press releases; Snowflake and "
             "AWS Marketplace listings; A-Team Insight; WatersTechnology.")

# ==================================================================== PAGE 2
s2 = page_frame(
    prs, layouts,
    "GoldenSource",
    "Market Position and Coverage Angle  |  July 2026",
    "Six of the category’s scaled independents have changed hands since 2022, most recently "
    "S&P Global’s EDM business to STG in January 2026. GoldenSource is one of three that is "
    "neither strategic-owned nor tied to a front-to-back platform.",
)

# --- zone A: competitive 2x2 (wide) + demand drivers (narrow)
cy = section(s2, LX, 1.56, 5.45, "COMPETITIVE LANDSCAPE")
matrix_2x2(s2, LX, cy, 5.45, 3.96 - cy, {
    (1, 0): ("INDEPENDENT ENTERPRISE MASTERS", [
        ("GoldenSource", {"bold": True, "color": TEAL, "font": HEAD}),
        (", Gresham (STG), NeoXam (Eurazeo). Cross-asset mastering sold on its own merits, "
         "not as a hook into proprietary data or an execution platform.", {})]),
    (1, 1): ("PLATFORM-TIED AND CAPTIVE SUITES", [
        ("Bloomberg, S&P Global, FactSet, LSEG, SimCorp (Deutsche Börse), Clearwater. Broad "
         "coverage, but the data model serves the owner’s feed or front-to-back stack.", {})]),
    (0, 0): ("POINT TOOLS AND INFRASTRUCTURE", [
        ("Snowflake, Databricks, Duco, Solidatus, Xceptor, Rimes. Neutral and increasingly "
         "capable, but each solves one slice; mastering is left to the client.", {})]),
    (0, 1): ("SINGLE-SOURCE FEEDS", [
        ("Vendor-native reference data delivery and in-house builds on one provider. Cheapest "
         "to start, hardest to reconcile once a second source arrives.", {})]),
}, "Breadth of mastered data (narrow to cross-asset)",
   "Vendor neutrality (low to high)")

cy = section(s2, 5.80, 1.56, 3.95, "DEMAND DRIVERS")
tf = textbox(s2, 5.82, cy, 3.91, 3.96 - cy)
write(tf, [
    lead("Data quality now gates AI. ",
         "InvestOps 2026 found 98% of firms concerned that poor data drives incorrect AI "
         "insights, moving governance from an operations budget to a board line item.",
         bullet=True),
    lead("Regulatory and cost pressure. ",
         "Rising vendor data costs, T+1, and widening entity, counterparty and ESG reporting "
         "obligations all reward a single mastered source.", bullet=True),
    lead("Cloud re-platforming. ",
         "Snowflake and Databricks made where the data lives a live decision again, and vendors "
         "that ship natively into those environments earn a second look.", bullet=True),
    lead("Vendor consolidation. ",
         "Front-to-back platforms squeeze point solutions but raise the premium on neutral "
         "mastering that spans them.", bullet=True),
], size=7.5, space_after=3.5)

# --- zone B: full-width consolidation timeline
cy = section(s2, LX, 4.06, FULL_W, "SELECTED PRECEDENT TRANSACTIONS")
timeline(s2, LX, cy, FULL_W, 5.16 - cy, [
    ("May-22", [("GoldenSource", True), ("Gemspring Capital", False), ("n.d.", False)], True),
    ("Apr-23", [("SimCorp", True), ("Deutsche Börse", False), ("€3.9B", False)], False),
    ("Jul-24", [("Gresham Technologies", True), ("STG, with Alveo", False), ("£141.9M", False)], False),
    ("Oct-24", [("EZOPS", True), ("NeoXam (Eurazeo)", False), ("n.d.", False)], False),
    ("Jan-25", [("Enfusion", True), ("Clearwater Analytics", False), ("$1.5B", False)], False),
    ("Jan-26", [("S&P EDM / thinkFolio", True), ("STG / Gresham", False), ("n.d.", False)], False),
], label_size=7)

# --- zone C: buyer universe (narrow) + coverage angle (wide)
cy = section(s2, LX, 5.26, 4.20, "INDICATIVE BUYER UNIVERSE")
lanes(s2, LX, cy, 4.20, BOT - cy, [
    ("DATA STRATEGICS", "S&P Global, LSEG, FactSet, Bloomberg, Deutsche Börse (SimCorp)"),
    ("PLATFORM CONSOLIDATORS", "Clearwater Analytics, SS&C, Broadridge, FIS"),
    ("CATEGORY SPONSORS", "STG (Alveo, Gresham, S&P EDM), Eurazeo (NeoXam), large-cap software funds"),
], gap=0.06)

cy = section(s2, 4.65, 5.26, 5.10, "COVERAGE ANGLE")
tf = textbox(s2, 4.67, cy, 5.06, BOT - cy)
write(tf, [
    lead("Hold period. ",
         "Gemspring entered in May 2022; at four-plus years and off a record 2025, the asset "
         "is in the window for a sponsor-to-sponsor or strategic process.", bullet=True),
    lead("Recent investment. ",
         "Scout, OMNI, V10 and the IBOR account for what the sponsor funded; headcount is up "
         "roughly 60% since close.", bullet=True),
    lead("Limited remaining supply. ",
         "The top-right quadrant above is down to three names, two of them already "
         "consolidated.", bullet=True),
    lead("Test in diligence. ",
         "Licence versus services mix; cloud ARR share against on-premise maintenance; net "
         "revenue retention; and how much of Scout has shipped.", bullet=True),
    lead("Relevant HL experience. ",
         "Rule 3 adviser to Gresham on its 2024 take-private by STG, and co-adviser to "
         "STG/Gresham on the 2026 S&P EDM and thinkFolio carve-out.",
         bullet=True, bullet_color=TEAL),
], size=7.5, space_after=2.5)

footnote(s2, "Sources: Houlihan Lokey transaction disclosures; S&P Global, Clearwater Analytics, "
             "Deutsche Börse and STG press releases; Businesswire; Finextra; A-Team Insight; "
             "InvestOps 2026. Values as disclosed; n.d. = not disclosed. Competitive positioning "
             "and buyer universe are HL analysis.")

prs.save(OUT)
print("wrote", OUT)

# ------------------------------- standalone: one slide holding only the panel
blk = Presentation(SRC)
delete_all_slides(blk)
bl = {l.name: l for l in blk.slide_master.slide_layouts}
bs = blk.slides.add_slide(bl["Blank"])
platform_panel(bs, 0.35, 2.65, 9.30, 2.20)
BLK_OUT = "GoldenSource_Platform_Block.pptx"
blk.save(BLK_OUT)
print("wrote", BLK_OUT)
