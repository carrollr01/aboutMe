"""For-profit conversion journey slide, built on the ACAMS master."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

SRC, OUT = "acams.pptx", "acams_journey.pptx"

# ---- ACAMS theme
NAVY   = RGBColor(0x18, 0x24, 0x44)   # accent2
BLUE   = RGBColor(0x12, 0x36, 0xD3)   # accent1
TEAL   = RGBColor(0x00, 0x84, 0x7E)   # accent5
MIDBL  = RGBColor(0x51, 0x75, 0xBD)   # accent6
GREY   = RGBColor(0x49, 0x49, 0x49)   # tx1
PALE   = RGBColor(0xD7, 0xED, 0xFF)   # lt2
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
PH_BG  = RGBColor(0xF4, 0xF5, 0xF7)
PH_LN  = RGBColor(0xD2, 0xD6, 0xDD)
RULE   = RGBColor(0xC9, 0xCD, 0xD4)
FONT   = "Figtree"

ROWS = [
    dict(name="Elemica", accent=BLUE, tint=RGBColor(0xE2,0xE7,0xFA),
         desc="Founded in 2002, developer of supply chain software by fourteen of "
              "the world’s leading chemicals companies",
         consortium=["DuPont", "Mitsubishi Gas Chemical", "Shell Chemicals"],
         first_pe=["Thoma Bravo"], buyout=["Eurazeo"],
         today="Eurazeo joined the cap table in 2019; company expanded offerings "
               "through a series of add-ons in recent years"),
    dict(name="Exostar", accent=MIDBL, tint=RGBColor(0xEA,0xEE,0xF7),
         desc="Founded as a joint venture among aerospace and defense companies to "
              "securely facilitate supply chain transactions",
         consortium=["BAE Systems", "Boeing", "Lockheed Martin", "Raytheon"],
         first_pe=["Thoma Bravo"], buyout=["ACP"],
         today="Since the acquisition in 2023, ACP has augmented the product "
               "offerings through acquisitions of ComplyUp and Robot Morning"),
    dict(name="GRESB", accent=TEAL, tint=RGBColor(0xE0,0xF0,0xEF),
         desc="Provider of sustainability standards established in 2009 by a "
              "consortium of institutional investors in collaboration with "
              "Maastricht University",
         consortium=["APG", "PGGM", "Maastricht University"],
         first_pe=["Summit Partners"], buyout=["General Atlantic"],
         today="General Atlantic acquired the business in 2019"),
    dict(name="Worldly", accent=NAVY, tint=RGBColor(0xDD,0xE2,0xED),
         desc="Sustainability and insights platform launched in 2019 as a "
              "public-benefit technology company",
         consortium=["Cascale"], first_pe=["Titan Grove"],
         buyout=["Buckhill", "Galvanize", "LFX Venture Partners",
                 "MissionPoint Partners", "Silversmith"],
         today="Raised $55m in a Series B round in 2022; Series A round in 2019 "
               "raised $4m"),
]
HEADS = ["Founding", "Founding Consortium¹", "First PE Investment",
         "Next Buyout", "Status Today"]

# ---- geometry (13.333 x 7.5)
LX, RX = 0.37, 12.96
CIRC_W = 1.02
BOX_X = LX + CIRC_W
BOX_W = RX - BOX_X
COLS = [2.55, 2.35, 2.20, 2.20, 2.27]        # sums to BOX_W = 11.57
HEAD_Y, HEAD_H = 1.30, 0.42
ROW_Y, ROW_H, ROW_GAP = 1.92, 1.00, 0.14


def no_line(sh):
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def rect(slide, x, y, w, h, fill=None, shape=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    return no_line(sh)


def tbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tf


def put(tf, text, size=8.5, color=GREY, bold=False, align=PP_ALIGN.LEFT,
        line_spacing=1.10):
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    p.space_after = Pt(0)
    r = p.add_run()
    r.text = text
    r.font.name, r.font.size, r.font.bold = FONT, Pt(size), bold
    r.font.color.rgb = color
    return tf


def logo_slot(slide, x, y, w, h, label):
    sh = rect(slide, x, y, w, h, PH_BG, MSO_SHAPE.ROUNDED_RECTANGLE)
    try:
        sh.adjustments[0] = 0.12
    except Exception:
        pass
    sh.line.fill.solid()
    sh.line.fill.fore_color.rgb = PH_LN
    sh.line.width = Pt(0.75)
    sh.name = "LOGO %s" % label
    tf = tbox(slide, x + 0.03, y, w - 0.06, h, anchor=MSO_ANCHOR.MIDDLE)
    put(tf, label, size=6.5, color=RGBColor(0x8A, 0x8F, 0x98),
        align=PP_ALIGN.CENTER, line_spacing=1.0)


def logo_grid(slide, x, y, w, h, labels):
    """Centred grid of logo placeholders inside a cell."""
    n = len(labels)
    per = 1 if n == 1 else (2 if n <= 4 else 3)
    rows = (n + per - 1) // per
    gx, gy = 0.07, 0.07
    lw = min((w - 0.16 - gx * (per - 1)) / per, 1.05)
    lh = min((h - 0.14 - gy * (rows - 1)) / rows, 0.30)
    total_h = rows * lh + (rows - 1) * gy
    y0 = y + (h - total_h) / 2
    i = 0
    for r in range(rows):
        in_row = min(per, n - i)
        total_w = in_row * lw + (in_row - 1) * gx
        x0 = x + (w - total_w) / 2
        for c in range(in_row):
            logo_slot(slide, x0 + c * (lw + gx), y0 + r * (lh + gy), lw, lh, labels[i])
            i += 1


prs = Presentation(SRC)
layouts = {l.name: l for l in prs.slide_masters[0].slide_layouts}
s = prs.slides.add_slide(layouts["Title Only"])

s.shapes.title.text_frame.text = ""
p = s.shapes.title.text_frame.paragraphs[0]
r = p.add_run()
r.text = ("For-Profit Conversion is a Reliable and Proven Journey for "
          "Successful Long-Term Value Creation")
r.font.name, r.font.size, r.font.bold = FONT, Pt(20), True
r.font.color.rgb = NAVY

# ---- column header chevrons
cx = BOX_X
for w, label in zip(COLS, HEADS):
    sh = rect(s, cx, HEAD_Y, w - 0.10, HEAD_H, WHITE, MSO_SHAPE.PENTAGON)
    try:
        sh.adjustments[0] = 0.14
    except Exception:
        pass
    sh.line.fill.solid()
    sh.line.fill.fore_color.rgb = BLUE
    sh.line.width = Pt(1.0)
    tf = tbox(s, cx + 0.06, HEAD_Y, w - 0.34, HEAD_H, anchor=MSO_ANCHOR.MIDDLE)
    put(tf, label, size=10, color=BLUE, bold=True, align=PP_ALIGN.CENTER,
        line_spacing=1.0)
    cx += w

# ---- rows
for i, row in enumerate(ROWS):
    ry = ROW_Y + i * (ROW_H + ROW_GAP)

    body = rect(s, BOX_X, ry, BOX_W, ROW_H, WHITE)
    body.line.fill.solid()
    body.line.fill.fore_color.rgb = row["accent"]
    body.line.width = Pt(0.75)

    circ = rect(s, LX, ry + (ROW_H - 0.86) / 2, 0.86, 0.86, row["tint"], MSO_SHAPE.OVAL)
    circ.line.fill.solid()
    circ.line.fill.fore_color.rgb = row["accent"]
    circ.line.width = Pt(1.0)
    circ.name = "LOGO %s" % row["name"]
    tf = tbox(s, LX + 0.05, ry + (ROW_H - 0.86) / 2, 0.76, 0.86, anchor=MSO_ANCHOR.MIDDLE)
    put(tf, row["name"], size=8, color=row["accent"], bold=True,
        align=PP_ALIGN.CENTER, line_spacing=1.0)

    cx = BOX_X
    tf = tbox(s, cx + 0.12, ry + 0.08, COLS[0] - 0.24, ROW_H - 0.16,
              anchor=MSO_ANCHOR.MIDDLE)
    put(tf, row["desc"], size=8.5)
    cx += COLS[0]

    for w, labels in ((COLS[1], row["consortium"]), (COLS[2], row["first_pe"]),
                      (COLS[3], row["buyout"])):
        rect(s, cx, ry + 0.10, 0.007, ROW_H - 0.20, RULE)      # column divider
        logo_grid(s, cx, ry, w, ROW_H, labels)
        cx += w

    rect(s, cx, ry + 0.10, 0.007, ROW_H - 0.20, RULE)
    tf = tbox(s, cx + 0.12, ry + 0.08, COLS[4] - 0.24, ROW_H - 0.16,
              anchor=MSO_ANCHOR.MIDDLE)
    put(tf, row["today"], size=8)

# ---- takeaway
ty = ROW_Y + len(ROWS) * (ROW_H + ROW_GAP) + 0.12
band = rect(s, LX, ty, RX - LX, 0.40, PALE, MSO_SHAPE.ROUNDED_RECTANGLE)
try:
    band.adjustments[0] = 0.18
except Exception:
    pass
tf = tbox(s, LX + 0.18, ty, RX - LX - 0.36, 0.40, anchor=MSO_ANCHOR.MIDDLE)
put(tf, "Increasing sponsor appetite in non-profit entities provides a proven "
        "playbook Sedex is well-positioned to replicate",
    size=11, color=NAVY, bold=True, align=PP_ALIGN.CENTER, line_spacing=1.0)

# ---- source note (drawn directly; the layout's footer is not a cloned placeholder)
tf = tbox(s, LX, 7.06, 10.71, 0.38)
for k, line in enumerate([
        "Source: Company websites, press releases and public filings.",
        "1.  Representative companies listed; does not contain the full list for "
        "Elemica and Exostar."]):
    pp = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
    pp.line_spacing = 1.05
    pp.space_after = Pt(0)
    rr = pp.add_run()
    rr.text = line
    rr.font.name, rr.font.size = FONT, Pt(7)
    rr.font.color.rgb = GREY

prs.save(OUT)
print("wrote", OUT, "| slide", len(prs.slides._sldIdLst))
